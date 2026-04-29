#!/usr/bin/env python3
"""
Tests for analyze_template.py and generate_doc.py.

Creates synthetic templates programmatically, runs the scripts against them,
and validates output structure and content.

Run: python3 evaluations/test_file_generation.py
"""
import json
import os
import subprocess
import sys
import tempfile
import unittest
from copy import copy

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ANALYZE = os.path.join(REPO_ROOT, ".claude", "scripts", "analyze_template.py")
GENERATE = os.path.join(REPO_ROOT, ".claude", "scripts", "generate_doc.py")


def run(script, *args):
    result = subprocess.run(
        [sys.executable, script, *args],
        capture_output=True, text=True
    )
    return result.returncode, result.stdout, result.stderr


# ─── Template factories ───────────────────────────────────────────────────────

def make_pptx(path, include_vo_shape=False, include_background=False):
    """
    Minimal storyboard .pptx with named layouts.
    include_vo_shape: adds a 'Voice Over Box' text box to the Storyboard Row layout.
    include_background: adds a background fill to the Knowledge Check layout.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_layouts[0].name = "Title Slide"
    prs.slide_layouts[1].name = "Storyboard Row"
    prs.slide_layouts[2].name = "Knowledge Check"

    if include_vo_shape:
        # Add a named text box to the Storyboard Row layout's slide master slide
        # (We add it to a test slide instead, as modifying layouts directly is complex)
        pass  # Tested via slide-level shape in test slides below

    if include_background:
        from pptx.oxml.ns import qn
        import lxml.etree as etree
        layout = prs.slide_layouts[2]
        bg = layout.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    prs.save(path)


def make_pptx_with_vo_shape(path):
    """Template with a named 'Voice Over Box' text box on each content slide."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_layouts[0].name = "Title Slide"
    prs.slide_layouts[1].name = "Storyboard Row"
    prs.slide_layouts[2].name = "Knowledge Check"

    # Add a test slide that has the VO shape — analyze picks up non-placeholder shapes
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(0), Inches(5), Inches(4), Inches(1))
    txBox.name = "Voice Over Box"
    txBox.text_frame.text = "VO placeholder"

    prs.save(path)


def make_pptx_with_multiple_prototypes(path):
    """Three Storyboard Row prototype slides, each with a uniquely named text box."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_layouts[0].name = "Title Slide"
    prs.slide_layouts[1].name = "Storyboard Row"
    prs.slide_layouts[2].name = "Knowledge Check"

    layout = prs.slide_layouts[1]
    for label in ("Box A", "Box B", "Box C"):
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0), Inches(4), Inches(3), Inches(1))
        tb.name = label
        tb.text_frame.text = f"Placeholder for {label}"

    prs.save(path)


def make_pptx_with_no_textframe_shape(path):
    """Prototype slide that includes a decorative shape element with no text body."""
    from pptx import Presentation
    from pptx.util import Inches
    import lxml.etree as etree

    prs = Presentation()
    prs.slide_layouts[0].name = "Title Slide"
    prs.slide_layouts[1].name = "Storyboard Row"

    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    tb = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1))
    tb.name = "Content Area"
    tb.text_frame.text = "Content placeholder"

    # Decorative line element — no <p:txBody>, so has_text_frame is False
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_xml = (
        f'<p:sp xmlns:p="{ns_p}" xmlns:a="{ns_a}">'
        "<p:nvSpPr>"
        '<p:cNvPr id="200" name="Decorative Line"/>'
        "<p:cNvSpPr/><p:nvPr/>"
        "</p:nvSpPr>"
        "<p:spPr>"
        '<a:xfrm><a:off x="0" y="4572000"/><a:ext cx="9144000" cy="50000"/></a:xfrm>'
        '<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        "</p:spPr>"
        "</p:sp>"
    )
    slide.shapes._spTree.append(etree.fromstring(sp_xml.encode("utf-8")))

    prs.save(path)


def make_xlsx_partial_columns(path):
    """xlsx with only 3 of the 6 standard alignment-matrix columns."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Alignment Matrix"
    for col, header in enumerate(
        ["Learning Objective", "Bloom's Level", "Instructional Strategy"], start=1
    ):
        ws.cell(row=1, column=col, value=header)
    wb.save(path)


def make_docx(path):
    from docx import Document
    from docx.shared import Pt
    from docx.oxml.ns import qn

    doc = Document()
    doc.add_heading("Workshop Title", level=1)
    doc.add_heading("Session Plan", level=2)

    tbl = doc.add_table(rows=2, cols=5)
    tbl.style = "Table Grid"
    headers = ["Time", "Activity", "Facilitator Action", "Anticipated Responses", "Notes"]
    for i, h in enumerate(headers):
        tbl.rows[0].cells[i].text = h
    sample = ["0:00", "Sample", "Sample action", "Sample response", ""]
    for i, v in enumerate(sample):
        tbl.rows[1].cells[i].text = v

    doc.add_heading("Troubleshooting", level=2)
    tbl2 = doc.add_table(rows=2, cols=2)
    tbl2.style = "Table Grid"
    tbl2.rows[0].cells[0].text = "Situation"
    tbl2.rows[0].cells[1].text = "Response"
    tbl2.rows[1].cells[0].text = "Sample"
    tbl2.rows[1].cells[1].text = "Sample"

    doc.save(path)


def make_xlsx(path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, Border, Side, Alignment

    wb = Workbook()
    ws = wb.active
    ws.title = "Alignment Matrix"
    headers = [
        "Learning Objective", "Bloom's Level", "Instructional Strategy",
        "Assessment Method", "Media Type", "Duration"
    ]
    thin = Side(style="thin")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for col, h in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True, name="Calibri", size=11)
        cell.border = border
        cell.alignment = Alignment(horizontal="center")
    wb.save(path)


# ─── Shared test data ─────────────────────────────────────────────────────────

# Mapping that routes narration to slide Notes (explicit, not hardcoded)
PPTX_MAPPING_NOTES = json.dumps({
    "slide_types": {
        "Title": "Title Slide",
        "Content": "Storyboard Row",
        "Scenario": "Storyboard Row",
        "Knowledge Check": "Knowledge Check",
    },
    "placeholder_map": {"0": "text"},
    "shape_map": {},
    "notes_fields": ["narration", "visual", "interaction", "dev_notes"],
})

# Mapping that routes narration to a named shape
PPTX_MAPPING_SHAPE = json.dumps({
    "slide_types": {
        "Title": "Title Slide",
        "Content": "Storyboard Row",
        "Scenario": "Storyboard Row",
        "Knowledge Check": "Knowledge Check",
    },
    "placeholder_map": {"0": "text"},
    "shape_map": {"Voice Over Box": "narration"},
    "notes_fields": ["visual", "dev_notes"],
})

PPTX_DATA = json.dumps({
    "title": "Delegate and Coach — Module 1",
    "version": "v1.0",
    "date": "2026-04-29",
    "slides": [
        {
            "num": 1, "type": "Title",
            "text": "Delegation and Coaching for New Managers",
            "narration": "Welcome to Module 1.",
            "visual": "Title slide with team image",
            "interaction": "None",
            "dev_notes": "Use brand header",
        },
        {
            "num": 2, "type": "Content",
            "text": "What does effective delegation look like?",
            "narration": "In this section we explore the four conditions for safe delegation.",
            "visual": "Split screen",
            "interaction": "Click to reveal",
            "dev_notes": "",
        },
        {
            "num": 3, "type": "Knowledge Check",
            "text": "Which sign shows a task is ready to delegate?",
            "narration": "",
            "visual": "MCQ layout",
            "interaction": "MCQ — 4 options",
            "dev_notes": "Correct: B.",
        },
    ],
})

XLSX_MAPPING = json.dumps({
    "column_map": {
        "objective": "Learning Objective",
        "bloom": "Bloom's Level",
        "strategy": "Instructional Strategy",
        "assessment": "Assessment Method",
        "media": "Media Type",
        "duration": "Duration",
    }
})

XLSX_DATA = json.dumps({
    "title": "Module 1 Alignment Matrix",
    "rows": [
        {
            "objective": "Given a delegation scenario, choose the appropriate task to delegate",
            "bloom": "Apply",
            "strategy": "Branching scenario",
            "assessment": "Scenario MCQ",
            "media": "Rise eLearning",
            "duration": "15 min",
        },
        {
            "objective": "Identify signs that a direct report is ready for increased autonomy",
            "bloom": "Analyze",
            "strategy": "Case study discussion",
            "assessment": "Observation checklist",
            "media": "ILT",
            "duration": "20 min",
        },
    ],
})

DOCX_MAPPING = json.dumps({
    "session_plan_table": 0,
    "column_map": {
        "time": "Time",
        "activity": "Activity",
        "action": "Facilitator Action",
        "responses": "Anticipated Responses",
        "notes": "Notes",
    },
    "troubleshooting_table": 1,
})

DOCX_DATA = json.dumps({
    "title": "New Manager Workshop — Day 1",
    "length": "4 hrs",
    "max_participants": 20,
    "format": "ILT",
    "materials": ["Slide deck", "Handout A"],
    "setup": "U-shape seating",
    "activities": [
        {
            "time": "0:00",
            "activity": "Welcome",
            "action": "Introduce yourself and the day's goals",
            "responses": "Participants share names and roles",
            "notes": "Keep energy high",
        },
        {
            "time": "0:15",
            "activity": "Delegation scenarios",
            "action": "Distribute scenario cards; small group discussion",
            "responses": "Groups identify delegation criteria",
            "notes": "Watch for 'everything is undelegatable' response",
        },
    ],
    "debrief": [
        "What surprised you about when to delegate?",
        "How will you apply this with your team this week?",
    ],
    "troubleshooting": [
        {
            "situation": "Group disengages during scenarios",
            "response": "Break into pairs; give one scenario each",
        }
    ],
})


# ─── analyze_template.py — pptx ───────────────────────────────────────────────

class TestAnalyzePptx(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.tmp.close()
        make_pptx(self.tmp.name)

    def tearDown(self):
        os.unlink(self.tmp.name)

    def test_returns_valid_json(self):
        rc, out, err = run(ANALYZE, "--file", self.tmp.name)
        self.assertEqual(rc, 0, err)
        data = json.loads(out)
        self.assertEqual(data["type"], "pptx")

    def test_layouts_present(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        names = [l["name"] for l in data["layouts"]]
        self.assertIn("Title Slide", names)
        self.assertIn("Storyboard Row", names)

    def test_each_layout_has_shapes_list(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        for layout in data["layouts"]:
            self.assertIn("shapes", layout)
            self.assertIsInstance(layout["shapes"], list)

    def test_layout_reports_has_background(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        for layout in data["layouts"]:
            self.assertIn("has_background", layout)

    def test_existing_slide_count(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        self.assertIn("existing_slide_count", data)
        self.assertIsInstance(data["existing_slide_count"], int)


class TestAnalyzePptxHints(unittest.TestCase):
    """Shape classification hints are surfaced correctly."""

    def setUp(self):
        self.tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.tmp.close()
        make_pptx_with_vo_shape(self.tmp.name)

    def tearDown(self):
        os.unlink(self.tmp.name)

    def test_vo_shape_detected_in_existing_slides(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        found_hints = []
        for slide in data.get("existing_slides", []):
            for shape in slide.get("non_placeholder_shapes", []):
                if "hint" in shape:
                    found_hints.append(shape["hint"])
        self.assertIn("candidate_narration", found_hints)

    def test_vo_shape_name_preserved(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        names = [
            s["name"]
            for slide in data.get("existing_slides", [])
            for s in slide.get("non_placeholder_shapes", [])
        ]
        self.assertIn("Voice Over Box", names)


class TestAnalyzeBackgroundDetection(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.tmp.close()
        make_pptx(self.tmp.name, include_background=True)

    def tearDown(self):
        os.unlink(self.tmp.name)

    def test_background_layout_flagged(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        kc = next((l for l in data["layouts"] if l["name"] == "Knowledge Check"), None)
        self.assertIsNotNone(kc)
        self.assertTrue(kc["has_background"])

    def test_plain_layout_not_flagged(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        data = json.loads(out)
        title_layout = next((l for l in data["layouts"] if l["name"] == "Title Slide"), None)
        self.assertIsNotNone(title_layout)
        self.assertFalse(title_layout["has_background"])


# ─── analyze_template.py — docx / xlsx / errors ───────────────────────────────

class TestAnalyzeDocx(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        self.tmp.close()
        make_docx(self.tmp.name)

    def tearDown(self):
        os.unlink(self.tmp.name)

    def test_returns_valid_json(self):
        rc, out, err = run(ANALYZE, "--file", self.tmp.name)
        self.assertEqual(rc, 0, err)
        self.assertEqual(json.loads(out)["type"], "docx")

    def test_detects_tables(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        self.assertGreaterEqual(len(json.loads(out)["tables"]), 1)

    def test_session_plan_headers_present(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        headers = json.loads(out)["tables"][0]["headers"]
        self.assertIn("Time", headers)
        self.assertIn("Facilitator Action", headers)

    def test_headings_detected(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        texts = [h["text"] for h in json.loads(out)["headings"]]
        self.assertIn("Session Plan", texts)


class TestAnalyzeXlsx(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.tmp.close()
        make_xlsx(self.tmp.name)

    def tearDown(self):
        os.unlink(self.tmp.name)

    def test_returns_valid_json(self):
        rc, out, err = run(ANALYZE, "--file", self.tmp.name)
        self.assertEqual(rc, 0, err)
        self.assertEqual(json.loads(out)["type"], "xlsx")

    def test_headers_detected(self):
        _, out, _ = run(ANALYZE, "--file", self.tmp.name)
        headers = json.loads(out)["sheets"][0]["headers"]
        self.assertIn("Learning Objective", headers)
        self.assertIn("Bloom's Level", headers)


class TestAnalyzeErrors(unittest.TestCase):
    def test_missing_file(self):
        rc, out, _ = run(ANALYZE, "--file", "/nonexistent/file.pptx")
        self.assertNotEqual(rc, 0)
        self.assertIn("error", json.loads(out))

    def test_unsupported_extension(self):
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"fake")
            name = f.name
        try:
            rc, out, _ = run(ANALYZE, "--file", name)
            self.assertNotEqual(rc, 0)
            self.assertIn("error", json.loads(out))
        finally:
            os.unlink(name)


# ─── generate_doc.py — storyboard (.pptx) ────────────────────────────────────

class TestGeneratePptxNotesRouting(unittest.TestCase):
    """Narration goes to slide Notes when notes_fields includes 'narration'."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f): os.unlink(f)

    def test_generates_file(self):
        rc, out, err = run(
            GENERATE,
            "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        self.assertTrue(os.path.exists(self.out.name))

    def test_slide_count_matches_data(self):
        from pptx import Presentation
        run(GENERATE, "--type", "storyboard", "--template", self.template.name,
            "--data", PPTX_DATA, "--mapping", PPTX_MAPPING_NOTES, "--out", self.out.name)
        prs = Presentation(self.out.name)
        self.assertEqual(len(prs.slides), len(json.loads(PPTX_DATA)["slides"]))

    def test_narration_in_notes(self):
        from pptx import Presentation
        run(GENERATE, "--type", "storyboard", "--template", self.template.name,
            "--data", PPTX_DATA, "--mapping", PPTX_MAPPING_NOTES, "--out", self.out.name)
        prs = Presentation(self.out.name)
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        self.assertIn("Welcome to Module 1", notes)

    def test_narration_not_in_notes_when_omitted_from_notes_fields(self):
        """When narration is routed to a shape, it should NOT also appear in notes."""
        from pptx import Presentation
        make_pptx_with_vo_shape(self.template.name)
        run(GENERATE, "--type", "storyboard", "--template", self.template.name,
            "--data", PPTX_DATA, "--mapping", PPTX_MAPPING_SHAPE, "--out", self.out.name)
        prs = Presentation(self.out.name)
        # Slide 0: narration = "Welcome to Module 1." — should NOT appear in notes
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        self.assertNotIn("Welcome to Module 1", notes)

    def test_knowledge_check_uses_correct_layout(self):
        from pptx import Presentation
        run(GENERATE, "--type", "storyboard", "--template", self.template.name,
            "--data", PPTX_DATA, "--mapping", PPTX_MAPPING_NOTES, "--out", self.out.name)
        prs = Presentation(self.out.name)
        self.assertEqual(prs.slides[2].slide_layout.name, "Knowledge Check")


class TestGeneratePptxShapeRouting(unittest.TestCase):
    """Narration goes to a named shape when shape_map routes it there."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx_with_vo_shape(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f): os.unlink(f)

    def test_generates_file(self):
        rc, out, err = run(
            GENERATE,
            "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", PPTX_MAPPING_SHAPE,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")

    def test_narration_in_named_shape(self):
        from pptx import Presentation
        run(GENERATE, "--type", "storyboard", "--template", self.template.name,
            "--data", PPTX_DATA, "--mapping", PPTX_MAPPING_SHAPE, "--out", self.out.name)
        prs = Presentation(self.out.name)
        # Template prototype slide is removed; output = [Title(0), Content(1), KnowledgeCheck(2)].
        # Content slide (index 1) uses Storyboard Row — the VO Box is copied from the prototype.
        slide = prs.slides[1]
        vo_shape = next(
            (s for s in slide.shapes if s.name == "Voice Over Box"), None
        )
        self.assertIsNotNone(vo_shape, "Voice Over Box shape should be present on Content slide")
        self.assertIn("In this section we explore", vo_shape.text_frame.text)

    def test_placeholder_map_alias_fields_key(self):
        """Legacy 'fields' key works as alias for 'placeholder_map'."""
        mapping_with_alias = json.loads(PPTX_MAPPING_NOTES)
        mapping_with_alias["fields"] = mapping_with_alias.pop("placeholder_map")
        rc, out, _ = run(
            GENERATE,
            "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", json.dumps(mapping_with_alias),
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0)


# ─── generate_doc.py — xlsx design inheritance ────────────────────────────────

class TestGenerateXlsx(unittest.TestCase):
    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.template.close()
        make_xlsx(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f): os.unlink(f)

    def test_generates_file(self):
        rc, out, err = run(
            GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name,
            "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        self.assertTrue(os.path.exists(self.out.name))

    def test_data_rows_written(self):
        from openpyxl import load_workbook
        run(GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name, "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING, "--out", self.out.name)
        ws = load_workbook(self.out.name).active
        rows = [r for r in ws.iter_rows(min_row=2, values_only=True) if any(v for v in r)]
        self.assertEqual(len(rows), 2)

    def test_bloom_cell_color_coded(self):
        from openpyxl import load_workbook
        run(GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name, "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING, "--out", self.out.name)
        ws = load_workbook(self.out.name).active
        bloom_col = next(
            (cell.column for cell in ws[1] if cell.value == "Bloom's Level"), None
        )
        self.assertIsNotNone(bloom_col)
        bloom_cell = ws.cell(row=2, column=bloom_col)
        self.assertNotEqual(bloom_cell.fill.fgColor.rgb, "00000000")

    def test_design_inheritance_border_copied(self):
        """Data cells inherit the border style from the header row."""
        from openpyxl import load_workbook
        run(GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name, "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING, "--out", self.out.name)
        wb = load_workbook(self.out.name)
        ws = wb.active
        header_cell = ws.cell(row=1, column=1)
        data_cell = ws.cell(row=2, column=1)
        # Both should have a border (header had thin borders; data row should inherit)
        self.assertIsNotNone(data_cell.border)
        self.assertEqual(
            data_cell.border.left.style,
            header_cell.border.left.style,
        )

    def test_design_inheritance_data_rows_not_bold(self):
        """Data rows should not inherit the header's bold flag."""
        from openpyxl import load_workbook
        run(GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name, "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING, "--out", self.out.name)
        ws = load_workbook(self.out.name).active
        data_cell = ws.cell(row=2, column=1)
        bold = data_cell.font.bold if data_cell.font else False
        self.assertFalse(bold)


# ─── generate_doc.py — docx ───────────────────────────────────────────────────

class TestGenerateDocx(unittest.TestCase):
    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        self.template.close()
        make_docx(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f): os.unlink(f)

    def test_generates_file(self):
        rc, out, err = run(
            GENERATE, "--type", "facilitator-guide",
            "--template", self.template.name,
            "--data", DOCX_DATA,
            "--mapping", DOCX_MAPPING,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        self.assertTrue(os.path.exists(self.out.name))

    def test_activity_rows_written(self):
        from docx import Document
        run(GENERATE, "--type", "facilitator-guide",
            "--template", self.template.name, "--data", DOCX_DATA,
            "--mapping", DOCX_MAPPING, "--out", self.out.name)
        tbl = Document(self.out.name).tables[0]
        self.assertEqual(len(tbl.rows), 3)  # header + 2 activities

    def test_activity_content_present(self):
        from docx import Document
        run(GENERATE, "--type", "facilitator-guide",
            "--template", self.template.name, "--data", DOCX_DATA,
            "--mapping", DOCX_MAPPING, "--out", self.out.name)
        tbl = Document(self.out.name).tables[0]
        all_text = " ".join(c.text for row in tbl.rows for c in row.cells)
        self.assertIn("Delegation scenarios", all_text)

    def test_troubleshooting_row_written(self):
        from docx import Document
        run(GENERATE, "--type", "facilitator-guide",
            "--template", self.template.name, "--data", DOCX_DATA,
            "--mapping", DOCX_MAPPING, "--out", self.out.name)
        tbl2 = Document(self.out.name).tables[1]
        all_text = " ".join(c.text for row in tbl2.rows for c in row.cells)
        self.assertIn("disengages", all_text)


# ─── Mapping persistence ──────────────────────────────────────────────────────

class TestMappingPersistence(unittest.TestCase):
    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.template.close()
        make_xlsx(self.template.name)
        self._real_mappings = os.path.join(REPO_ROOT, "templates", "mappings.json")
        self._backup = None
        if os.path.exists(self._real_mappings):
            with open(self._real_mappings) as f:
                self._backup = f.read()

    def tearDown(self):
        os.unlink(self.template.name)
        if self._backup is not None:
            with open(self._real_mappings, "w") as f:
                f.write(self._backup)
        elif os.path.exists(self._real_mappings):
            with open(self._real_mappings, "w") as f:
                f.write("{}\n")

    def test_save_mapping(self):
        rc, out, err = run(
            GENERATE, "--save-mapping",
            "--template", self.template.name,
            "--mapping", XLSX_MAPPING,
        )
        self.assertEqual(rc, 0, err)
        self.assertIn("saved", json.loads(out))

    def test_saved_mapping_keyed_by_basename(self):
        run(GENERATE, "--save-mapping",
            "--template", self.template.name, "--mapping", XLSX_MAPPING)
        with open(self._real_mappings) as f:
            mappings = json.load(f)
        self.assertIn(os.path.basename(self.template.name), mappings)

    def test_missing_mapping_returns_exit_2(self):
        tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        tmp.close()
        make_xlsx(tmp.name)
        try:
            rc, out, _ = run(
                GENERATE, "--type", "alignment-matrix",
                "--template", tmp.name, "--data", XLSX_DATA,
            )
            self.assertEqual(rc, 2)
            self.assertEqual(json.loads(out)["error"], "no_mapping")
        finally:
            os.unlink(tmp.name)


# ─── High-stress edge cases ───────────────────────────────────────────────────
#
# Cases 1–10 exercise failure modes not covered by the baseline suite:
# unicode/XML-unsafe chars, empty data, unknown slide types, duplicate field
# routing, all-empty fields, multiple prototypes, non-text-frame shapes,
# missing xlsx columns, large slide count, and nonexistent shape names.


class TestStressPptx(unittest.TestCase):
    """Cases 1, 2, 3, 5, 9, 10 — plain pptx template (no prototype slides)."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f):
                os.unlink(f)

    # ── Case 1: Unicode and XML-unsafe characters in all text fields ──────────

    def test_unicode_and_xml_unsafe_chars(self):
        """Unicode, emoji, smart quotes, XML-special chars must not crash or corrupt."""
        data = json.dumps({
            "title": "Test <Course> & ‘Module’ \U0001f3af",
            "version": "v1.0",
            "date": "2026-04-29",
            "slides": [{
                "num": 1, "type": "Title",
                "text": "H\xe9llo <World> & \"Friends\"",
                "narration": "Smart “quotes” and em—dash\nnewline here",
                "visual": "Ampersand & angle <brackets>",
                "interaction": "None",
                "dev_notes": "Special chars: <>&\"'",
            }],
        })
        rc, out, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        from pptx import Presentation
        prs = Presentation(self.out.name)
        self.assertEqual(len(prs.slides), 1)

    # ── Case 2: Empty slides list ─────────────────────────────────────────────

    def test_empty_slides_generates_valid_file(self):
        """Zero slides in data → valid openable pptx with 0 slides, exit 0."""
        data = json.dumps({
            "title": "Empty", "version": "v1.0", "date": "2026-04-29", "slides": []
        })
        rc, out, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(self.out.name).slides), 0)

    # ── Case 3: Unknown slide type not in mapping ─────────────────────────────

    def test_unknown_slide_type_falls_to_default_layout(self):
        """A slide type absent from slide_types falls to default layout, no crash."""
        data = json.dumps({
            "title": "Test", "version": "v1.0", "date": "2026-04-29",
            "slides": [{"num": 1, "type": "Quiz", "text": "What is delegation?",
                        "narration": "Unknown type slide", "visual": "", "interaction": "",
                        "dev_notes": ""}],
        })
        rc, _, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(self.out.name).slides), 1)

    # ── Case 5: All data fields empty for a slide ─────────────────────────────

    def test_all_fields_empty_generates_slide_without_crash(self):
        """A slide with every field empty still generates cleanly; no crash or KeyError."""
        data = json.dumps({
            "title": "Test", "version": "v1.0", "date": "2026-04-29",
            "slides": [{
                "num": 1, "type": "Content",
                "text": "", "narration": "", "visual": "",
                "interaction": "", "dev_notes": "",
            }],
        })
        rc, _, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(self.out.name).slides), 1)

    # ── Case 9: Large storyboard — 50 slides ─────────────────────────────────

    def test_large_storyboard_50_slides(self):
        """50-slide generation completes, output opens, and slide count is exact."""
        types = ["Title", "Content", "Scenario", "Knowledge Check"]
        slides = [
            {
                "num": i + 1,
                "type": types[i % len(types)],
                "text": f"Slide {i + 1} on-screen text",
                "narration": f"Narration for slide {i + 1}.",
                "visual": f"Visual description {i + 1}",
                "interaction": "None",
                "dev_notes": f"Dev note {i + 1}",
            }
            for i in range(50)
        ]
        data = json.dumps({
            "title": "Large Storyboard", "version": "v1.0", "date": "2026-04-29",
            "slides": slides,
        })
        rc, _, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(self.out.name).slides), 50)

    # ── Case 10: shape_map references a shape that exists in no prototype ──────

    def test_shape_map_nonexistent_shape_silently_skipped(self):
        """shape_map entry for a shape not in any prototype is silently skipped."""
        mapping = json.dumps({
            "slide_types": {
                "Title": "Title Slide", "Content": "Storyboard Row",
                "Scenario": "Storyboard Row", "Knowledge Check": "Knowledge Check",
            },
            "placeholder_map": {"0": "text"},
            "shape_map": {"Nonexistent Box": "narration"},
            "notes_fields": [],
        })
        rc, _, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", mapping,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}")
        from pptx import Presentation
        self.assertEqual(len(Presentation(self.out.name).slides), 3)


class TestStressDuplicateRouting(unittest.TestCase):
    """Case 4 — field appears in both shape_map and notes_fields; shape_map must win."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx_with_vo_shape(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f):
                os.unlink(f)

    def test_shape_map_wins_over_notes_fields(self):
        """When a field is in both shape_map and notes_fields, it goes to the shape only."""
        mapping = json.dumps({
            "slide_types": {
                "Title": "Title Slide", "Content": "Storyboard Row",
                "Scenario": "Storyboard Row", "Knowledge Check": "Knowledge Check",
            },
            "placeholder_map": {"0": "text"},
            "shape_map": {"Voice Over Box": "narration"},
            "notes_fields": ["narration", "visual"],  # narration in BOTH
        })
        run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", mapping,
            "--out", self.out.name,
        )
        from pptx import Presentation
        prs = Presentation(self.out.name)
        content_slide = prs.slides[1]  # Content slide — has VO Box from prototype

        vo_shape = next(
            (s for s in content_slide.shapes if s.name == "Voice Over Box"), None
        )
        self.assertIsNotNone(vo_shape, "Voice Over Box must be present on Content slide")
        self.assertIn("In this section we explore", vo_shape.text_frame.text)

        notes_text = content_slide.notes_slide.notes_text_frame.text
        self.assertNotIn(
            "In this section we explore", notes_text,
            "Narration must not appear in notes when routed to a shape",
        )
        # visual is only in notes_fields (not shape_map) — it should be in notes
        self.assertIn("Split screen", notes_text)


class TestStressMultiplePrototypes(unittest.TestCase):
    """Case 6 — three prototype slides for the same layout; only first is used."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx_with_multiple_prototypes(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f):
                os.unlink(f)

    def test_only_first_prototype_shapes_copied(self):
        """When a layout has multiple prototype slides only the first one's shapes copy."""
        run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", PPTX_DATA,
            "--mapping", PPTX_MAPPING_NOTES,
            "--out", self.out.name,
        )
        from pptx import Presentation
        prs = Presentation(self.out.name)
        self.assertEqual(len(prs.slides), 3)  # 3 template slides removed; 3 generated

        content_slide = prs.slides[1]  # Storyboard Row layout
        shape_names = [s.name for s in content_slide.shapes]
        self.assertIn("Box A", shape_names, "First prototype shape must be copied")
        self.assertNotIn("Box B", shape_names, "Second prototype must not bleed in")
        self.assertNotIn("Box C", shape_names, "Third prototype must not bleed in")


class TestStressNonTextFrameShape(unittest.TestCase):
    """Case 7 — prototype contains a non-text-frame shape (decorative element)."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.template.close()
        make_pptx_with_no_textframe_shape(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f):
                os.unlink(f)

    def test_non_text_frame_shape_copies_without_crash(self):
        """Decorative shapes without a text body deepcopy cleanly; output is valid."""
        data = json.dumps({
            "title": "Test", "version": "v1.0", "date": "2026-04-29",
            "slides": [{"num": 1, "type": "Content", "text": "On-screen text",
                        "narration": "Narration here", "visual": "", "interaction": "",
                        "dev_notes": ""}],
        })
        mapping = json.dumps({
            "slide_types": {"Content": "Storyboard Row"},
            "placeholder_map": {},
            "shape_map": {},
            "notes_fields": ["narration"],
        })
        rc, out, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", mapping,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")
        from pptx import Presentation
        prs = Presentation(self.out.name)
        self.assertEqual(len(prs.slides), 1)
        shape_names = [s.name for s in prs.slides[0].shapes]
        self.assertIn("Content Area", shape_names, "Text-frame shape must survive deepcopy")

    def test_populate_skips_non_text_frame_shape_in_shape_map(self):
        """shape_map entry pointing to a non-text-frame shape is silently skipped."""
        data = json.dumps({
            "title": "Test", "version": "v1.0", "date": "2026-04-29",
            "slides": [{"num": 1, "type": "Content", "text": "", "narration": "Audio here",
                        "visual": "", "interaction": "", "dev_notes": ""}],
        })
        mapping = json.dumps({
            "slide_types": {"Content": "Storyboard Row"},
            "placeholder_map": {},
            "shape_map": {"Decorative Line": "narration"},  # no text frame — must skip
            "notes_fields": [],
        })
        rc, _, err = run(
            GENERATE, "--type", "storyboard",
            "--template", self.template.name,
            "--data", data,
            "--mapping", mapping,
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}")


class TestStressXlsxPartialColumns(unittest.TestCase):
    """Case 8 — column_map references columns absent from the xlsx template."""

    def setUp(self):
        self.template = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.template.close()
        make_xlsx_partial_columns(self.template.name)
        self.out = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self.out.close()

    def tearDown(self):
        for f in [self.template.name, self.out.name]:
            if os.path.exists(f):
                os.unlink(f)

    def test_missing_columns_silently_skipped(self):
        """Columns referenced in column_map but absent in the template are ignored."""
        rc, out, err = run(
            GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name,
            "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING,  # maps all 6 columns; template has 3
            "--out", self.out.name,
        )
        self.assertEqual(rc, 0, f"stderr: {err}\nstdout: {out}")

    def test_present_columns_are_written(self):
        """The 3 columns that do exist in the template receive data."""
        from openpyxl import load_workbook
        run(
            GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name,
            "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING,
            "--out", self.out.name,
        )
        ws = load_workbook(self.out.name).active
        col_headers = {cell.value for cell in ws[1] if cell.value}
        self.assertIn("Learning Objective", col_headers)
        obj_col = next(c.column for c in ws[1] if c.value == "Learning Objective")
        self.assertIsNotNone(ws.cell(row=2, column=obj_col).value)

    def test_absent_columns_not_written(self):
        """Columns present in column_map but absent from the template have no data."""
        from openpyxl import load_workbook
        run(
            GENERATE, "--type", "alignment-matrix",
            "--template", self.template.name,
            "--data", XLSX_DATA,
            "--mapping", XLSX_MAPPING,
            "--out", self.out.name,
        )
        ws = load_workbook(self.out.name).active
        col_headers = {cell.value for cell in ws[1] if cell.value}
        self.assertNotIn("Assessment Method", col_headers)
        self.assertNotIn("Media Type", col_headers)
        self.assertNotIn("Duration", col_headers)


if __name__ == "__main__":
    unittest.main(verbosity=2)
