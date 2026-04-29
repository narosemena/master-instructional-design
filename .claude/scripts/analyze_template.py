#!/usr/bin/env python3
"""
Analyze a template file and output its structure as JSON.

Usage:
  python3 analyze_template.py --file PATH

Output (stdout): JSON describing the template's structural elements.

  .pptx  — all shapes per layout (placeholders + free text boxes), semantic
            hints for VO/dev-notes/visual containers, background presence,
            existing slide count and per-slide preview
  .docx  — paragraph styles used, heading hierarchy, table column headers
  .xlsx  — sheet names, header row values, row count per sheet

Shape hints in pptx output:
  candidate_narration   — shape name contains VO/narration/script/audio signal
  candidate_dev_notes   — shape name contains dev/production note signal
  candidate_visual      — shape name contains visual/art direction signal
  candidate_interaction — shape name contains interaction/interactivity signal

These hints surface during excavation so field destinations are confirmed,
never assumed.
"""
import argparse
import json
import os
import sys

# Keyword sets used to classify shape names as likely content containers
_VO_SIGNALS      = frozenset(["voice over", "vo", "narration", "script", "audio", "speaker"])
_DEV_SIGNALS     = frozenset(["dev note", "developer", "dev notes", "production note", "for dev"])
_VISUAL_SIGNALS  = frozenset(["visual", "image desc", "art direction", "graphic desc", "illustration"])
_INTERACT_SIGNALS = frozenset(["interaction", "interaction type", "interactivity", "user action"])


def _shape_hint(name):
    """Return a semantic hint if the shape name signals a specific content type."""
    n = name.lower()
    for sig in _VO_SIGNALS:
        if sig in n:
            return "candidate_narration"
    for sig in _DEV_SIGNALS:
        if sig in n:
            return "candidate_dev_notes"
    for sig in _VISUAL_SIGNALS:
        if sig in n:
            return "candidate_visual"
    for sig in _INTERACT_SIGNALS:
        if sig in n:
            return "candidate_interaction"
    return None


def _layout_has_background(layout):
    """True if the layout has an explicit non-transparent background fill."""
    try:
        from pptx.oxml.ns import qn
        elem = layout._element
        cSld = elem.find(qn('p:cSld'))
        if cSld is None:
            return False
        bg = cSld.find(qn('p:bg'))
        if bg is None:
            return False
        bgPr = bg.find(qn('p:bgPr'))
        if bgPr is None:
            return False
        for fill_tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            if bgPr.find(qn(fill_tag)) is not None:
                return True
        return False
    except Exception:
        return False


def analyze_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)

    layouts = []
    for layout in prs.slide_layouts:
        shapes = []
        for shape in layout.shapes:
            info = {
                "name": shape.name,
                "has_text_frame": shape.has_text_frame,
                "is_placeholder": shape.is_placeholder,
            }
            if shape.is_placeholder:
                info["placeholder_idx"] = shape.placeholder_format.idx
                info["placeholder_type"] = (
                    str(shape.placeholder_format.type).split(".")[-1]
                )
            hint = _shape_hint(shape.name)
            if hint:
                info["hint"] = hint
            shapes.append(info)

        layouts.append({
            "name": layout.name,
            "has_background": _layout_has_background(layout),
            "shapes": shapes,
        })

    existing_slides = []
    for i, slide in enumerate(prs.slides):
        non_ph_shapes = []
        for shape in slide.shapes:
            if not shape.is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                entry = {"name": shape.name}
                hint = _shape_hint(shape.name)
                if hint:
                    entry["hint"] = hint
                if text:
                    entry["preview"] = text[:80]
                non_ph_shapes.append(entry)

        existing_slides.append({
            "num": i + 1,
            "layout": slide.slide_layout.name,
            "non_placeholder_shapes": non_ph_shapes,
        })

    return {
        "type": "pptx",
        "existing_slide_count": len(prs.slides),
        "existing_slides": existing_slides,
        "layouts": layouts,
    }


def analyze_docx(path):
    from docx import Document

    doc = Document(path)

    styles_used = sorted({p.style.name for p in doc.paragraphs if p.text.strip()})
    headings = [
        {"style": p.style.name, "text": p.text}
        for p in doc.paragraphs
        if p.style.name.startswith("Heading") and p.text.strip()
    ]

    tables = []
    for i, tbl in enumerate(doc.tables):
        headers = [c.text.strip() for c in tbl.rows[0].cells] if tbl.rows else []
        tables.append({
            "table_num": i + 1,
            "headers": headers,
            "row_count": len(tbl.rows),
            "col_count": len(tbl.columns),
        })

    return {
        "type": "docx",
        "styles_used": styles_used,
        "headings": headings,
        "tables": tables,
    }


def analyze_xlsx(path):
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True)
    sheets = []
    for name in wb.sheetnames:
        ws = wb[name]
        headers = []
        for row in ws.iter_rows(min_row=1, max_row=1, values_only=True):
            headers = [str(h) if h is not None else "" for h in row]
        sheets.append({
            "name": name,
            "headers": headers,
            "row_count": ws.max_row or 0,
        })
    wb.close()
    return {"type": "xlsx", "sheets": sheets}


def main():
    parser = argparse.ArgumentParser(
        description="Analyze a template file and output its structure as JSON"
    )
    parser.add_argument("--file", required=True, help="Path to template file")
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(json.dumps({"error": f"File not found: {args.file}"}))
        sys.exit(1)

    ext = os.path.splitext(args.file)[1].lower()
    try:
        if ext == ".pptx":
            result = analyze_pptx(args.file)
        elif ext in (".docx", ".doc"):
            result = analyze_docx(args.file)
        elif ext in (".xlsx", ".xls"):
            result = analyze_xlsx(args.file)
        else:
            print(json.dumps({"error": f"Unsupported file type: {ext}"}))
            sys.exit(1)
        print(json.dumps(result, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
