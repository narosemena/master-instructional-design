#!/usr/bin/env python3
"""
Generate a populated document from a user template using a confirmed field mapping.

Usage:
  # Generate (uses saved mapping from templates/mappings.json if available):
  python3 generate_doc.py --type TYPE --template PATH --data JSON [--out PATH]

  # Save a confirmed mapping without generating:
  python3 generate_doc.py --save-mapping --template PATH --mapping JSON

  # Generate with inline mapping and save it for future sessions:
  python3 generate_doc.py --type TYPE --template PATH --data JSON \
    --mapping JSON --save-mapping [--out PATH]

Supported types and output formats:
  storyboard        → .pptx
  facilitator-guide → .docx
  alignment-matrix  → .xlsx

─────────────────────────────────────────────────────────────────────────────
DATA SCHEMAS
─────────────────────────────────────────────────────────────────────────────

storyboard:
  {
    "title": "Course — Module 1",
    "version": "v1.0",
    "date": "YYYY-MM-DD",
    "slides": [
      {
        "num": 1,
        "type": "Title | Content | Scenario | Knowledge Check",
        "text": "On-screen text",
        "narration": "Narration / voice-over script",
        "visual": "Visual description for developer",
        "interaction": "None | Click to reveal | Branch choice | MCQ",
        "dev_notes": "Developer production notes"
      }
    ]
  }

facilitator-guide:
  {
    "title": "Workshop Name",
    "length": "4 hrs",
    "max_participants": 20,
    "format": "ILT | VILT",
    "materials": ["Slide deck", "Handout A"],
    "setup": "Room / platform setup notes",
    "activities": [
      {
        "time": "0:00",
        "activity": "Opening",
        "action": "Facilitator action description",
        "responses": "Anticipated participant responses",
        "notes": "Internal notes"
      }
    ],
    "debrief": ["Question 1", "Question 2"],
    "troubleshooting": [
      {"situation": "Group disengages", "response": "Break into pairs"}
    ]
  }

alignment-matrix:
  {
    "title": "Module Alignment Matrix",
    "rows": [
      {
        "objective": "Full objective text",
        "bloom": "Remember | Understand | Apply | Analyze | Evaluate | Create",
        "strategy": "Instructional strategy",
        "assessment": "Assessment method",
        "media": "Media type",
        "duration": "X min"
      }
    ]
  }

─────────────────────────────────────────────────────────────────────────────
MAPPING SCHEMAS  (stored in templates/mappings.json, keyed by template basename)
─────────────────────────────────────────────────────────────────────────────

storyboard (.pptx):
  {
    "slide_types": {
      "Title": "Title Slide",        ← slide type → layout name in template
      "Content": "Storyboard Row",
      "Scenario": "Storyboard Row",
      "Knowledge Check": "Knowledge Check"
    },
    "placeholder_map": {
      "0": "text",                   ← placeholder idx → data field
      "1": "visual"
    },
    "shape_map": {
      "Voice Over Box": "narration", ← shape name → data field
      "Dev Notes": "dev_notes"
    },
    "notes_fields": ["interaction"]  ← fields appended to slide Notes
  }

  Field routing priority: placeholder_map → shape_map → notes_fields.
  Any field not in any map is silently skipped.
  Destinations are confirmed during excavation — nothing is hardcoded.

  Alias: "fields" is accepted as a synonym for "placeholder_map" for
  compatibility with hand-authored mappings.

facilitator-guide (.docx):
  {
    "session_plan_table": 0,         ← 0-based table index for session plan
    "column_map": {
      "time": "Time",                ← data field → column header in template
      "activity": "Activity",
      "action": "Facilitator Action",
      "responses": "Anticipated Responses",
      "notes": "Notes"
    },
    "troubleshooting_table": 1       ← optional: 0-based index of troubleshooting table
  }

alignment-matrix (.xlsx):
  {
    "column_map": {
      "objective": "Learning Objective",   ← data field → column header in template
      "bloom": "Bloom's Level",
      "strategy": "Instructional Strategy",
      "assessment": "Assessment Method",
      "media": "Media Type",
      "duration": "Duration"
    }
  }
"""
import argparse
import json
import os
import sys


MAPPINGS_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
    "templates",
    "mappings.json",
)

TYPE_EXT = {
    "storyboard": "pptx",
    "facilitator-guide": "docx",
    "alignment-matrix": "xlsx",
}

BLOOM_COLORS = {
    "Remember": "DBEAFE",
    "Understand": "BBF7D0",
    "Apply": "FEF08A",
    "Analyze": "FED7AA",
    "Evaluate": "FECACA",
    "Create": "E9D5FF",
}


# ─── Mapping persistence ──────────────────────────────────────────────────────

def load_mappings():
    if os.path.exists(MAPPINGS_PATH):
        with open(MAPPINGS_PATH) as f:
            return json.load(f)
    return {}


def save_mapping(template_basename, mapping):
    mappings = load_mappings()
    mappings[template_basename] = mapping
    os.makedirs(os.path.dirname(MAPPINGS_PATH), exist_ok=True)
    with open(MAPPINGS_PATH, "w") as f:
        json.dump(mappings, f, indent=2)
    print(json.dumps({"saved": template_basename, "path": MAPPINGS_PATH}))


def resolve_mapping(template_path, mapping_arg):
    """Return (mapping_dict, source_description)."""
    if mapping_arg:
        if mapping_arg.strip().startswith("{"):
            return json.loads(mapping_arg), "inline"
        with open(mapping_arg) as f:
            return json.load(f), mapping_arg

    basename = os.path.basename(template_path)
    mappings = load_mappings()
    if basename in mappings:
        return mappings[basename], f"templates/mappings.json ({basename})"

    return None, None


# ─── Silent design inheritance helpers ───────────────────────────────────────

def _copy_cell_format(source_cell, target_cell, skip_fill=False):
    """Copy non-content formatting from source to target cell (silent; skips on error)."""
    from copy import copy
    try:
        if source_cell.font:
            f = copy(source_cell.font)
            f.bold = False  # Data rows are never bold; headers may be
            target_cell.font = f
        if source_cell.border:
            target_cell.border = copy(source_cell.border)
        if source_cell.alignment:
            target_cell.alignment = copy(source_cell.alignment)
        if source_cell.number_format:
            target_cell.number_format = source_cell.number_format
        if not skip_fill and source_cell.fill and source_cell.fill.fill_type not in (None, "none"):
            target_cell.fill = copy(source_cell.fill)
    except Exception:
        pass  # Design inheritance is always best-effort


def _copy_docx_row_format(source_row, target_row):
    """Copy paragraph and cell formatting from a docx table row (silent; skips on error)."""
    from copy import deepcopy
    try:
        for src_cell, tgt_cell in zip(source_row.cells, target_row.cells):
            for src_para, tgt_para in zip(src_cell.paragraphs, tgt_cell.paragraphs):
                if src_para._pPr is not None:
                    tgt_para._p.get_or_add_pPr()
                    tgt_para._pPr._p = deepcopy(src_para._pPr._p)
    except Exception:
        pass


def _layout_has_background(layout):
    """True if the layout has an explicit non-transparent background fill."""
    try:
        from pptx.oxml.ns import qn
        elem = layout._element
        cSld = elem.find(qn('p:cSld'))
        if cSld is None:
            return False
        bg = cSld.find(qn('p:bg'))
        if bg is None:
            return False
        bgPr = bg.find(qn('p:bgPr'))
        if bgPr is None:
            return False
        for fill_tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
            if bgPr.find(qn(fill_tag)) is not None:
                return True
        return False
    except Exception:
        return False


# ─── XLSX generator ───────────────────────────────────────────────────────────

def generate_xlsx(template_path, data, mapping, output_path):
    from openpyxl import load_workbook
    from openpyxl.styles import PatternFill

    wb = load_workbook(template_path)
    ws = wb.active

    column_map = mapping.get("column_map", {})
    fields = ["objective", "bloom", "strategy", "assessment", "media", "duration"]

    # Find header row (first non-empty row within first 10)
    header_row_idx = 1
    for i, row in enumerate(ws.iter_rows(min_row=1, max_row=10, values_only=True), start=1):
        if any(v for v in row):
            header_row_idx = i
            break

    header_cells = list(ws.iter_rows(
        min_row=header_row_idx, max_row=header_row_idx
    ))[0]
    header_to_col = {
        cell.value: cell.column
        for cell in header_cells
        if cell.value is not None
    }

    col_indices = {
        field: header_to_col[column_map[field]]
        for field in fields
        if column_map.get(field) in header_to_col
    }

    # Bloom's column index for selective fill handling
    bloom_col = col_indices.get("bloom")

    start_row = header_row_idx + 1
    for i, row_data in enumerate(data.get("rows", [])):
        row_idx = start_row + i
        for field, col_idx in col_indices.items():
            cell = ws.cell(row=row_idx, column=col_idx)
            cell.value = row_data.get(field, "")

            # Silent design inheritance: copy header row formatting
            header_cell = ws.cell(row=header_row_idx, column=col_idx)
            is_bloom = (col_idx == bloom_col)
            _copy_cell_format(header_cell, cell, skip_fill=is_bloom)

            # Bloom's color coding applied after format copy
            if is_bloom:
                color = BLOOM_COLORS.get(row_data.get("bloom", ""))
                if color:
                    cell.fill = PatternFill(
                        start_color=color, end_color=color, fill_type="solid"
                    )

    wb.save(output_path)


# ─── DOCX generator ───────────────────────────────────────────────────────────

def generate_docx(template_path, data, mapping, output_path):
    from docx import Document

    doc = Document(template_path)

    column_map = mapping.get("column_map", {
        "time": "Time",
        "activity": "Activity",
        "action": "Facilitator Action",
        "responses": "Anticipated Responses",
        "notes": "Notes",
    })
    session_plan_idx = mapping.get("session_plan_table", 0)
    troubleshooting_idx = mapping.get("troubleshooting_table", None)

    # Replace title in the first Title/Heading 1 paragraph
    for para in doc.paragraphs:
        if para.style.name in ("Title", "Heading 1") and para.text.strip():
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = data.get("title", "Facilitator Guide")
            break

    # Populate session plan table
    if doc.tables and session_plan_idx < len(doc.tables):
        tbl = doc.tables[session_plan_idx]
        if tbl.rows:
            col_by_header = {
                cell.text.strip(): i for i, cell in enumerate(tbl.rows[0].cells)
            }

        act_fields = ["time", "activity", "action", "responses", "notes"]
        col_order = [
            col_by_header.get(column_map.get(f, ""), -1) for f in act_fields
        ]

        # Capture the last existing data row for format inheritance before clearing
        format_source_row = tbl.rows[-1] if len(tbl.rows) > 1 else tbl.rows[0]

        # Remove existing data rows (keep header)
        for row in tbl.rows[1:]:
            tbl._tbl.remove(row._tr)

        for activity in data.get("activities", []):
            row = tbl.add_row()
            _copy_docx_row_format(format_source_row, row)
            for field, col_idx in zip(act_fields, col_order):
                if 0 <= col_idx < len(row.cells):
                    row.cells[col_idx].text = str(activity.get(field, ""))

    # Populate troubleshooting table if mapped
    if troubleshooting_idx is not None and doc.tables and troubleshooting_idx < len(doc.tables):
        tbl = doc.tables[troubleshooting_idx]
        format_source_row = tbl.rows[-1] if len(tbl.rows) > 1 else tbl.rows[0]
        for row in tbl.rows[1:]:
            tbl._tbl.remove(row._tr)
        for item in data.get("troubleshooting", []):
            row = tbl.add_row()
            _copy_docx_row_format(format_source_row, row)
            if len(row.cells) >= 2:
                row.cells[0].text = item.get("situation", "")
                row.cells[1].text = item.get("response", "")

    doc.save(output_path)


# ─── PPTX generator ───────────────────────────────────────────────────────────

def _populate_slide(slide, slide_data, mapping):
    """
    Populate a slide's content using the three-tier field routing:
      1. placeholder_map  — {placeholder_idx_str: data_field}
      2. shape_map        — {shape_name: data_field}
      3. notes_fields     — [data_field, ...]  → appended to slide Notes

    Nothing is hardcoded. All destinations are confirmed during excavation and
    stored in the mapping. Fields absent from all three maps are silently skipped.
    """
    # Support "fields" as a legacy alias for "placeholder_map"
    placeholder_map = mapping.get("placeholder_map") or mapping.get("fields", {})
    shape_map = mapping.get("shape_map", {})
    notes_fields = mapping.get("notes_fields", [])

    # Track fields routed to slide content; these are excluded from notes (priority rule).
    routed_fields = set()

    # 1. Placeholder population
    for ph in slide.placeholders:
        idx_str = str(ph.placeholder_format.idx)
        field = placeholder_map.get(idx_str)
        if field and slide_data.get(field):
            try:
                ph.text = str(slide_data[field])
                routed_fields.add(field)
            except Exception:
                pass

    # 2. Named shape population
    shape_by_name = {
        s.name: s for s in slide.shapes
        if s.has_text_frame and not s.is_placeholder
    }
    for shape_name, field in shape_map.items():
        if slide_data.get(field):
            shape = shape_by_name.get(shape_name)
            if shape:
                try:
                    shape.text_frame.text = str(slide_data[field])
                    routed_fields.add(field)
                except Exception:
                    pass

    # 3. Notes field population — fields already routed to slide content are skipped
    if notes_fields:
        parts = [
            str(slide_data[f]) for f in notes_fields
            if slide_data.get(f) and f not in routed_fields
        ]
        if parts:
            slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)


def generate_pptx(template_path, data, mapping, output_path):
    from pptx import Presentation
    from copy import deepcopy

    prs = Presentation(template_path)

    # Collect non-placeholder shapes from template slides as per-layout prototypes
    proto_shapes = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        if name not in proto_shapes:
            non_ph = [s for s in slide.shapes if not s.is_placeholder]
            if non_ph:
                proto_shapes[name] = non_ph

    template_slide_count = len(prs.slides)

    slide_types = mapping.get("slide_types", {})
    layout_by_name = {layout.name: layout for layout in prs.slide_layouts}
    default_layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]

    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type", "Content")
        layout_name = slide_types.get(slide_type)
        layout = layout_by_name.get(layout_name, default_layout)

        slide = prs.slides.add_slide(layout)

        # Copy prototype non-placeholder shapes to preserve template design elements
        if layout.name in proto_shapes:
            spTree = slide.shapes._spTree
            max_id = max(
                (int(el.get('id', 0)) for el in spTree.iter() if 'cNvPr' in el.tag and el.get('id')),
                default=0,
            )
            for proto_shape in proto_shapes[layout.name]:
                try:
                    new_elem = deepcopy(proto_shape._element)
                    max_id += 1
                    for el in new_elem.iter():
                        if 'cNvPr' in el.tag:
                            el.set('id', str(max_id))
                            break
                    spTree.append(new_elem)
                except Exception:
                    pass

        _populate_slide(slide, slide_data, mapping)

    # Remove template prototype slides — output contains only generated content
    sldIdLst = prs.slides._sldIdLst
    for _ in range(template_slide_count):
        if len(sldIdLst):
            sldIdLst.remove(sldIdLst[0])

    prs.save(output_path)


# ─── Output path helper ───────────────────────────────────────────────────────

def build_output_path(doc_type, title, _template_path):
    ext = TYPE_EXT.get(doc_type, "")
    slug = title.lower().replace(" ", "-").replace("—", "").replace("/", "-")
    slug = "".join(c for c in slug if c.isalnum() or c == "-").strip("-")[:60]
    out_dir = os.path.join(
        os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
        "generated",
    )
    os.makedirs(out_dir, exist_ok=True)
    return os.path.join(out_dir, f"{slug}.{ext}")


# ─── Entry point ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Populate a user template with structured L&D content"
    )
    parser.add_argument("--type", choices=list(TYPE_EXT.keys()))
    parser.add_argument("--template", required=True)
    parser.add_argument("--data")
    parser.add_argument("--mapping")
    parser.add_argument("--save-mapping", action="store_true")
    parser.add_argument("--out")
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print(json.dumps({"error": f"Template not found: {args.template}"}))
        sys.exit(1)

    # Save-mapping-only mode
    if args.save_mapping and not args.data:
        if not args.mapping:
            print(json.dumps({"error": "--mapping required with --save-mapping"}))
            sys.exit(1)
        mapping = (
            json.loads(args.mapping) if args.mapping.strip().startswith("{")
            else json.load(open(args.mapping))
        )
        save_mapping(os.path.basename(args.template), mapping)
        return

    if not args.type:
        print(json.dumps({"error": "--type required for generation"}))
        sys.exit(1)

    if not args.data:
        print(json.dumps({"error": "--data required for generation"}))
        sys.exit(1)

    mapping, mapping_source = resolve_mapping(args.template, args.mapping)
    if mapping is None:
        print(json.dumps({
            "error": "no_mapping",
            "message": (
                f"No mapping found for '{os.path.basename(args.template)}'. "
                "Run analyze_template.py to inspect the template structure, "
                "confirm field destinations via excavation, then pass "
                "--mapping JSON --save-mapping."
            ),
        }))
        sys.exit(2)

    if args.save_mapping and args.mapping:
        save_mapping(os.path.basename(args.template), mapping)

    data = json.loads(args.data)
    title = data.get("title", args.type)
    output_path = args.out or build_output_path(args.type, title, args.template)

    try:
        if args.type == "storyboard":
            generate_pptx(args.template, data, mapping, output_path)
        elif args.type == "facilitator-guide":
            generate_docx(args.template, data, mapping, output_path)
        elif args.type == "alignment-matrix":
            generate_xlsx(args.template, data, mapping, output_path)

        print(json.dumps({
            "output": output_path,
            "type": args.type,
            "mapping_source": mapping_source,
            "items": len(data.get(
                "slides", data.get("rows", data.get("activities", []))
            )),
        }))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)


if __name__ == "__main__":
    main()
